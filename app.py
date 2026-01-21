import streamlit as st
from streamlit_option_menu import option_menu
from pypdf import PdfReader, PdfWriter
import pdfplumber
import ollama
import io
import zipfile
import base64
import time
import random
import re
import json
import pandas as pd
import fitz  # PyMuPDF
import difflib
import os
import subprocess
import spacy
import numpy as np
from pptx import Presentation
from streamlit_agraph import agraph, Node, Edge, Config
from streamlit_timeline import timeline
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from docx import Document
from pdf2docx import Converter as DocxConverter

# --- 1. PAGE CONFIGURATION ---
st.set_page_config(page_title="Journey Through Pages", layout="wide", page_icon="📜")

# --- 2. PREMIUM UI ARCHITECTURE (ENHANCED FOR PUBLIC HOSTING) ---
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&family=Poppins:wght@400;500;600;700&display=swap');
    
    /* Global App Style - Single Solid Color */
    .stApp { 
        background: linear-gradient(180deg, #f0f4f8 0%, #e2e8f0 100%);
        font-family: 'Inter', sans-serif;
    }
    
    /* Sidebar Enhancement */
    .css-1d391kg { background: rgba(255, 255, 255, 0.1) !important; backdrop-filter: blur(30px) !important; border-right: 1px solid rgba(255, 255, 255, 0.2) !important; }
    
    /* Navigation Menu - Premium Styling */
    .stSidebar [data-testid="stSidebar"] {
        background: rgba(255, 255, 255, 0.08) !important;
        backdrop-filter: blur(30px) !important;
    }
    
    /* Animations */
    @keyframes fadeInUp { 
        from { opacity: 0; transform: translate3d(0, 30px, 0); } 
        to { opacity: 1; transform: translate3d(0, 0, 0); } 
    }
    @keyframes shimmer {
        0% { background-position: -1000px 0; }
        100% { background-position: 1000px 0; }
    }
    .animate-enter { animation: fadeInUp 0.6s ease-out; }
    
    /* Enhanced Glassmorphism Cards - Refined & Subtle */
    .glass-card { 
        background: rgba(255, 255, 255, 0.85); 
        backdrop-filter: blur(25px) saturate(180%);
        -webkit-backdrop-filter: blur(25px) saturate(180%);
        border: 1px solid rgba(255, 255, 255, 0.75); 
        box-shadow: 0 8px 32px 0 rgba(31, 38, 135, 0.12),
                    0 2px 8px 0 rgba(0, 0, 0, 0.05),
                    inset 0 1px 0 0 rgba(255, 255, 255, 0.9); 
        border-radius: 24px; 
        padding: 28px; 
        margin-bottom: 24px; 
        transition: all 0.4s cubic-bezier(0.4, 0, 0.2, 1);
        height: 100%;
        position: relative;
        overflow: hidden;
    }
    .glass-card::before {
        content: '';
        position: absolute;
        top: 0;
        left: -100%;
        width: 100%;
        height: 100%;
        background: linear-gradient(90deg, transparent, rgba(255, 255, 255, 0.4), transparent);
        transition: left 0.5s;
    }
    .glass-card:hover::before {
        left: 100%;
    }
    .glass-card:hover { 
        transform: translateY(-8px) scale(1.02); 
        box-shadow: 0 20px 60px 0 rgba(37, 99, 235, 0.2),
                    0 4px 12px 0 rgba(0, 0, 0, 0.08),
                    inset 0 1px 0 0 rgba(255, 255, 255, 0.95); 
        border-color: rgba(37, 99, 235, 0.4); 
    }
    
    /* Typography - Enhanced */
    h1 { 
        background: linear-gradient(135deg, #667eea 0%, #764ba2 50%, #f093fb 100%); 
        -webkit-background-clip: text; 
        -webkit-text-fill-color: transparent; 
        background-clip: text;
        font-weight: 800; 
        letter-spacing: -1px; 
        font-family: 'Poppins', sans-serif;
        text-shadow: 0 2px 20px rgba(102, 126, 234, 0.3);
    }
    h2, h3 { 
        color: #1e293b; 
        font-weight: 700; 
        margin-top: 0;
        font-family: 'Poppins', sans-serif;
    }
    
    /* Uniform Home Page Cards */
    .home-card-content {
        min-height: 200px;
        display: flex;
        flex-direction: column;
        position: relative;
        z-index: 1;
    }
    .home-card-content ul { 
        padding-left: 24px; 
        color: #475569; 
        font-size: 0.95rem; 
        line-height: 1.8; 
        list-style: none;
    }
    .home-card-content ul li::before {
        content: "▸";
        color: #667eea;
        font-weight: bold;
        display: inline-block;
        width: 1em;
        margin-left: -1em;
    }
    
    /* Premium Buttons */
    div.stButton > button { 
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
        color: white; 
        border: none; 
        border-radius: 12px; 
        padding: 0.75rem 1.5rem; 
        font-weight: 600; 
        font-family: 'Inter', sans-serif;
        width: 100%; 
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1); 
        box-shadow: 0 4px 14px 0 rgba(102, 126, 234, 0.4),
                    0 2px 4px 0 rgba(0, 0, 0, 0.1);
        position: relative;
        overflow: hidden;
    }
    div.stButton > button::before {
        content: '';
        position: absolute;
        top: 50%;
        left: 50%;
        width: 0;
        height: 0;
        border-radius: 50%;
        background: rgba(255, 255, 255, 0.3);
        transform: translate(-50%, -50%);
        transition: width 0.6s, height 0.6s;
    }
    div.stButton > button:hover::before {
        width: 300px;
        height: 300px;
    }
    div.stButton > button:hover { 
        transform: translateY(-3px) scale(1.02); 
        box-shadow: 0 8px 25px 0 rgba(102, 126, 234, 0.5),
                    0 4px 8px 0 rgba(0, 0, 0, 0.15);
    }
    div.stButton > button:active {
        transform: translateY(-1px) scale(0.98);
    }
    
    /* Enhanced Tabs Styling */
    .stTabs [data-baseweb="tab-list"] { 
        gap: 12px; 
        background-color: transparent; 
        padding: 8px 0;
    }
    .stTabs [data-baseweb="tab"] {
        height: 48px;
        background: rgba(255, 255, 255, 0.7);
        backdrop-filter: blur(10px);
        border-radius: 12px;
        padding: 0 24px;
        border: 1px solid rgba(255, 255, 255, 0.5);
        font-weight: 500;
        transition: all 0.3s;
        color: #475569;
    }
    .stTabs [data-baseweb="tab"]:hover {
        background: rgba(255, 255, 255, 0.85);
        transform: translateY(-2px);
    }
    .stTabs [aria-selected="true"] {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        box-shadow: 0 4px 14px 0 rgba(102, 126, 234, 0.4);
        border-color: transparent;
    }
    
    /* Premium Selectbox Styling */
    div[data-baseweb="select"] > div {
        background: rgba(255, 255, 255, 0.9) !important;
        backdrop-filter: blur(10px) !important;
        border: 1px solid rgba(102, 126, 234, 0.2) !important;
        border-radius: 12px !important;
        padding: 8px 12px !important;
        transition: all 0.3s !important;
        box-shadow: 0 2px 8px rgba(0, 0, 0, 0.05) !important;
    }
    div[data-baseweb="select"] > div:hover {
        border-color: rgba(102, 126, 234, 0.4) !important;
        box-shadow: 0 4px 12px rgba(102, 126, 234, 0.15) !important;
    }
    div[data-baseweb="select"] > div:focus-within {
        border-color: #667eea !important;
        box-shadow: 0 0 0 3px rgba(102, 126, 234, 0.1) !important;
    }
    
    /* Premium File Uploader Styling */
    .stFileUploader > div {
        background: rgba(255, 255, 255, 0.85) !important;
        backdrop-filter: blur(15px) !important;
        border: 2px dashed rgba(102, 126, 234, 0.3) !important;
        border-radius: 16px !important;
        padding: 24px !important;
        transition: all 0.3s !important;
    }
    .stFileUploader > div:hover {
        border-color: rgba(102, 126, 234, 0.5) !important;
        background: rgba(255, 255, 255, 0.95) !important;
        transform: translateY(-2px);
        box-shadow: 0 8px 24px rgba(102, 126, 234, 0.15) !important;
    }
    
    /* Text Input Styling */
    .stTextInput > div > div > input {
        background: rgba(255, 255, 255, 0.9) !important;
        backdrop-filter: blur(10px) !important;
        border: 1px solid rgba(102, 126, 234, 0.2) !important;
        border-radius: 12px !important;
        padding: 10px 16px !important;
        transition: all 0.3s !important;
    }
    .stTextInput > div > div > input:focus {
        border-color: #667eea !important;
        box-shadow: 0 0 0 3px rgba(102, 126, 234, 0.1) !important;
        background: rgba(255, 255, 255, 0.95) !important;
    }
    
    /* Number Input Styling */
    .stNumberInput > div > div > input {
        background: rgba(255, 255, 255, 0.9) !important;
        backdrop-filter: blur(10px) !important;
        border: 1px solid rgba(102, 126, 234, 0.2) !important;
        border-radius: 12px !important;
        padding: 10px 16px !important;
    }
    
    /* Radio Button Styling */
    .stRadio > div {
        background: rgba(255, 255, 255, 0.7) !important;
        backdrop-filter: blur(10px) !important;
        border-radius: 12px !important;
        padding: 12px !important;
        border: 1px solid rgba(102, 126, 234, 0.2) !important;
    }
    
    /* Checkbox Styling */
    .stCheckbox > label {
        background: rgba(255, 255, 255, 0.7) !important;
        backdrop-filter: blur(10px) !important;
        border-radius: 10px !important;
        padding: 8px 12px !important;
        border: 1px solid rgba(102, 126, 234, 0.2) !important;
        transition: all 0.3s !important;
    }
    .stCheckbox > label:hover {
        background: rgba(255, 255, 255, 0.85) !important;
        border-color: rgba(102, 126, 234, 0.4) !important;
    }
    
    /* Container Border Enhancement */
    div[data-testid="stVerticalBlockBorderWrapper"] {
        background: rgba(255, 255, 255, 0.8) !important;
        backdrop-filter: blur(20px) !important;
        border: 1px solid rgba(255, 255, 255, 0.6) !important;
        border-radius: 20px !important;
        padding: 24px !important;
        box-shadow: 0 8px 32px rgba(31, 38, 135, 0.1) !important;
    }
    
    /* Helper Classes - Enhanced */
    .confidence-badge { 
        display: inline-flex; 
        align-items: center; 
        padding: 6px 16px; 
        border-radius: 24px; 
        font-size: 0.85rem; 
        font-weight: 700; 
        margin-bottom: 12px;
        background: linear-gradient(135deg, rgba(102, 126, 234, 0.2), rgba(118, 75, 162, 0.2));
        backdrop-filter: blur(10px);
        border: 1px solid rgba(255, 255, 255, 0.5);
        box-shadow: 0 2px 8px rgba(0, 0, 0, 0.1);
    }
    .book-item { 
        padding: 14px 18px; 
        background: rgba(255, 255, 255, 0.9); 
        backdrop-filter: blur(10px);
        border-radius: 12px; 
        margin-bottom: 10px; 
        border-left: 4px solid #667eea; 
        font-size: 0.9rem; 
        color: #1e293b; 
        display: flex; 
        justify-content: space-between; 
        align-items: center; 
        box-shadow: 0 2px 8px rgba(0,0,0,0.08);
        transition: all 0.3s;
    }
    .book-item:hover {
        transform: translateX(4px);
        box-shadow: 0 4px 12px rgba(102, 126, 234, 0.2);
    }
    .diff-add { 
        background: linear-gradient(135deg, rgba(34, 197, 94, 0.15), rgba(22, 101, 52, 0.1)); 
        color: #166534; 
        padding: 14px 18px; 
        border-radius: 12px; 
        margin: 8px 0; 
        border-left: 4px solid #22c55e; 
        backdrop-filter: blur(10px);
        box-shadow: 0 2px 8px rgba(34, 197, 94, 0.1);
    }
    .diff-rem { 
        background: linear-gradient(135deg, rgba(239, 68, 68, 0.15), rgba(153, 27, 27, 0.1)); 
        color: #991b1b; 
        padding: 14px 18px; 
        border-radius: 12px; 
        margin: 8px 0; 
        border-left: 4px solid #ef4444; 
        backdrop-filter: blur(10px);
        box-shadow: 0 2px 8px rgba(239, 68, 68, 0.1);
    }
    
    /* Info/Success/Warning/Error Boxes */
    .stAlert {
        border-radius: 16px !important;
        backdrop-filter: blur(10px) !important;
        border: 1px solid rgba(255, 255, 255, 0.5) !important;
    }
    
    /* Download Button Enhancement */
    .stDownloadButton > button {
        background: linear-gradient(135deg, #10b981 0%, #059669 100%) !important;
        border-radius: 12px !important;
        box-shadow: 0 4px 14px rgba(16, 185, 129, 0.4) !important;
    }
    
    /* Expander Styling */
    .streamlit-expanderHeader {
        background: rgba(255, 255, 255, 0.7) !important;
        backdrop-filter: blur(10px) !important;
        border-radius: 12px !important;
        border: 1px solid rgba(102, 126, 234, 0.2) !important;
    }
    
    /* Top Navigation Bar */
    .top-nav {
        background: rgba(255, 255, 255, 0.9);
        backdrop-filter: blur(20px);
        border-bottom: 1px solid rgba(102, 126, 234, 0.2);
        padding: 12px 24px;
        margin: -1rem -1rem 2rem -1rem;
        display: flex;
        align-items: center;
        justify-content: space-between;
        box-shadow: 0 2px 10px rgba(0, 0, 0, 0.05);
        position: sticky;
        top: 0;
        z-index: 100;
    }
    .nav-buttons {
        display: flex;
        gap: 8px;
        flex-wrap: wrap;
    }
    .nav-btn {
        padding: 8px 16px;
        border-radius: 10px;
        background: rgba(102, 126, 234, 0.1);
        color: #667eea;
        text-decoration: none;
        font-weight: 500;
        font-size: 0.9rem;
        transition: all 0.3s;
        border: 1px solid rgba(102, 126, 234, 0.2);
        cursor: pointer;
    }
    .nav-btn:hover {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(102, 126, 234, 0.3);
    }
    .nav-btn.active {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
    }
    .home-btn {
        display: inline-flex;
        align-items: center;
        gap: 8px;
        padding: 10px 20px;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border-radius: 12px;
        text-decoration: none;
        font-weight: 600;
        transition: all 0.3s;
        box-shadow: 0 4px 12px rgba(102, 126, 234, 0.3);
    }
    .home-btn:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 20px rgba(102, 126, 234, 0.4);
    }
    
    /* Hide Streamlit Branding */
    #MainMenu { visibility: hidden; }
    footer { visibility: hidden; }
    header { visibility: hidden; }
    
</style>
""", unsafe_allow_html=True)

# --- 3. UI HELPER COMPONENTS ---
def ui_spacer(height=20): st.markdown(f"<div style='height: {height}px;'></div>", unsafe_allow_html=True)
def ui_skeleton(): 
    return st.empty().markdown("""
    <div style="
        background: linear-gradient(90deg, #667eea 0%, #764ba2 50%, #667eea 100%);
        background-size: 200% 100%;
        height: 100px;
        border-radius: 16px;
        animation: shimmer 2s infinite;
        box-shadow: 0 4px 20px rgba(102, 126, 234, 0.3);
        display: flex;
        align-items: center;
        justify-content: center;
        color: white;
        font-weight: 600;
        font-size: 1.1rem;
    ">
        <span style="animation: pulse 1.5s infinite;">⏳ Processing...</span>
    </div>
    <style>
        @keyframes shimmer {
            0% { background-position: -200% 0; }
            100% { background-position: 200% 0; }
        }
        @keyframes pulse {
            0%, 100% { opacity: 1; }
            50% { opacity: 0.5; }
        }
    </style>
    """, unsafe_allow_html=True)
def ui_confidence_badge(score="High"):
    c = {"High": "#dcfce7", "Medium": "#fef9c3", "Low": "#fee2e2"}
    st.markdown(f"<div class='confidence-badge' style='background-color: {c.get(score)}; color: #374151;'>⚡ AI Confidence: {score}</div>", unsafe_allow_html=True)

# --- 4. SESSION STATE ---
if 'page_selection' not in st.session_state: st.session_state.page_selection = 0 
if 'highlights' not in st.session_state: st.session_state.highlights = []
if 'preview_pdf_bytes' not in st.session_state: st.session_state.preview_pdf_bytes = None
if 'bookshelf' not in st.session_state: st.session_state.bookshelf = []

def set_page(index): st.session_state.page_selection = index

# --- TOP NAVIGATION BAR ---
def render_top_nav(current_page):
    pages = ["Home", "Viewer", "Converter", "Editor", "AI Analyst", "Compare", "About"]
    page_icons = ["🏠", "👀", "🔄", "✂️", "🧠", "⚖️", "ℹ️"]
    page_indices = {"Home": 0, "Viewer": 1, "Converter": 2, "Editor": 3, "AI Analyst": 4, "Compare": 5, "About": 6}
    
    nav_html = f"""
    <div class="top-nav">
        <div style="display: flex; align-items: center; gap: 12px;">
            <span style="font-size: 1.5rem;">📜</span>
            <span style="font-weight: 700; font-size: 1.1rem; color: #1e293b;">Journey Through Pages</span>
        </div>
        <div class="nav-buttons">
    """
    
    for page, icon in zip(pages, page_icons):
        is_active = "active" if page == current_page else ""
        nav_html += f"""
            <a href="javascript:void(0)" onclick="window.parent.postMessage({{type: 'streamlit:setComponentValue', value: {page_indices[page]}}}, '*')" 
               class="nav-btn {is_active}">
                {icon} {page}
            </a>
        """
    
    nav_html += """
        </div>
    </div>
    """
    return nav_html

# Removed render_home_button - no longer needed with top navigation

# --- 5. BOOKSHELF LOGIC ---
def add_to_bookshelf(filename, file_bytes):
    for book in st.session_state.bookshelf:
        if book['name'] == filename: return 
    st.session_state.bookshelf.insert(0, {'name': filename, 'data': file_bytes}) 
    if len(st.session_state.bookshelf) > 5: st.session_state.bookshelf.pop()

def get_file_from_shelf(filename):
    for book in st.session_state.bookshelf:
        if book['name'] == filename: return io.BytesIO(book['data'])
    return None

def remove_from_shelf(filename):
    st.session_state.bookshelf = [b for b in st.session_state.bookshelf if b['name'] != filename]

# --- 6. CORE FUNCTIONS ---
def get_daily_quote():
    return random.choice([
        "“Data is the new oil.” — Clive Humby", 
        "“Simplicity is the ultimate sophistication.” — Leonardo da Vinci", 
        "“Technology is best when it brings people together.” — Matt Mullenweg", 
        "“The best way to predict the future is to invent it.” — Alan Kay"
    ])

def extract_text_with_references(file):
    pages_data = []
    with pdfplumber.open(file) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text: pages_data.append({"page": i + 1, "text": text})
    return pages_data

def text_to_audio_mac(text):
    output_file = "podcast.wav"
    safe_text = text.replace('"', '').replace("'", "").replace("\n", " ")[:5000]
    try:
        subprocess.run(["say", "-o", output_file, "--data-format=LEI16@44100", safe_text], check=True)
        if os.path.exists(output_file):
            with open(output_file, "rb") as f: return f.read()
    except: return None

# --- ML FUNCTIONS ---
def extract_timeline_data(text):
    prompt = f"Extract key events. Return STRICT JSON: {{'events': [{{'start_date': {{'year': 'YYYY', 'month': 'MM', 'day': 'DD'}}, 'text': {{'headline': 'Title', 'text': 'Desc'}}}}]}}. Text: {text[:4000]}"
    try:
        res = ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': prompt}])
        match = re.search(r'\{.*\}', res['message']['content'], re.DOTALL)
        return json.loads(match.group()) if match else None
    except: return None

def generate_ppt_from_text(text):
    prompt = f"Create 5-slide PPT. Format: TITLE: [T]\nCONTENT: [B1, B2]. Text: {text[:4000]}"
    try:
        res = ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': prompt}])
        prs = Presentation(); prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = "Generated Presentation"
        for t, c in re.findall(r'TITLE: (.*?)\nCONTENT: (.*?)(?=\nTITLE:|$)', res['message']['content'], re.DOTALL):
            s = prs.slides.add_slide(prs.slide_layouts[1]); s.shapes.title.text = t.strip()
            tf = s.shapes.placeholders[1].text_frame
            for b in c.split(','): tf.add_paragraph().text = b.strip()
        out = io.BytesIO(); prs.save(out); return out
    except: return None

def get_embedding(text):
    try: return ollama.embeddings(model='llama3.2', prompt=text)["embedding"]
    except: return []

def semantic_search_pdf(file_bytes, query):
    doc = fitz.open(stream=file_bytes, filetype="pdf"); qv = get_embedding(query)
    if not qv: return None, 0
    best_score = -1; best_page = None; best_sent = None
    for p in doc:
        for s in p.get_text().replace('\n',' ').split('. '):
            if len(s) > 20:
                v = get_embedding(s)
                if v:
                    score = np.dot(qv, v) / (np.linalg.norm(qv) * np.linalg.norm(v))
                    if score > best_score: best_score = score; best_sent = s; best_page = p
    if best_score > 0.35 and best_page:
        for i in best_page.search_for(best_sent[:50]): best_page.add_highlight_annot(i).set_colors(stroke=(0,1,0)); best_page.add_highlight_annot(i).update()
    out = io.BytesIO(); doc.save(out); return out, best_score

def run_audit(text):
    try: return ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': f"Audit Risk Score 0-100. 3 Risks. Format: SCORE: X. RISKS: -. Text: {text[:4000]}"}])['message']['content']
    except: return "Audit Failed."

def analyze_contradictions(text):
    try: return ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': f"Find contradictions. Text: {text[:4000]}"}])['message']['content']
    except: return "Analysis Failed."

def extract_ledger_data(text):
    try:
        res = ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': f"Extract invoice JSON: {{'vendor': '', 'total': '', 'items': []}}. Text: {text[:3000]}"}])
        match = re.search(r'\{.*\}', res['message']['content'], re.DOTALL)
        return json.loads(match.group()) if match else None
    except: return None

def generate_flashcards(text):
    try:
        res = ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': f"5 Flashcards 'Term|Definition'. No headers. Text: {text[:3000]}"}])
        return [{'Question': l.split('|')[0], 'Answer': l.split('|')[1]} for l in res['message']['content'].split('\n') if '|' in l]
    except: return []

# --- STANDARD UTILS ---
def apply_redaction(file_bytes, search_text):
    doc = fitz.open(stream=file_bytes, filetype="pdf"); count = 0
    for page in doc:
        for inst in page.search_for(search_text): page.add_redact_annot(inst, fill=(0, 0, 0)); page.apply_redactions(); count += 1
    out = io.BytesIO(); doc.save(out); doc.close(); return out, count

def apply_highlights_to_pdf(file_bytes, highlights_list):
    doc = fitz.open(stream=file_bytes, filetype="pdf"); 
    for h in highlights_list:
        for page in doc:
            for inst in page.search_for(h['text']): page.add_highlight_annot(inst).update()
    out = io.BytesIO(); doc.save(out); doc.close(); return out, []

def detect_pii(text):
    return list(set(re.findall(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', text) + re.findall(r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b', text)))

def df_to_pdf(df):
    buffer = io.BytesIO(); doc = SimpleDocTemplate(buffer, pagesize=letter); elements = []
    data = [df.columns.to_list()] + df.values.tolist()
    table = Table(data)
    table.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.grey),('TEXTCOLOR',(0,0),(-1,0),colors.whitesmoke),('ALIGN',(0,0),(-1,-1),'CENTER'),('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),('BOTTOMPADDING',(0,0),(-1,0),12),('BACKGROUND',(0,1),(-1,-1),colors.beige),('GRID',(0,0),(-1,-1),1,colors.black)]))
    elements.append(table); doc.build(elements); return buffer

def pdf_to_docx(f):
    out = io.BytesIO(); open("t.pdf", "wb").write(f.getvalue())
    cv = DocxConverter("t.pdf"); cv.convert("t.docx"); cv.close()
    out.write(open("t.docx", "rb").read()); return out

def pdf_to_excel(f):
    # Simplified placeholder for PDF to Excel
    return io.BytesIO(b"Excel conversion placeholder")

def images_to_pdf(imgs):
    out = io.BytesIO(); i1 = Image.open(imgs[0]).convert('RGB'); i1.save(out, save_all=True, append_images=[Image.open(i).convert('RGB') for i in imgs[1:]], format='PDF'); return out

def word_to_pdf(d):
    doc = Document(d); out = io.BytesIO(); c = canvas.Canvas(out, pagesize=letter); t = c.beginText(40, 750); t.setFont("Helvetica", 12)
    for p in doc.paragraphs: t.textLine(p.text[:90]); c.drawText(t) if t.getY() > 40 else (c.showPage(), t.setTextOrigin(40, 750))
    c.save(); return out

def split_pdf(file):
    r = PdfReader(file); z = io.BytesIO()
    with zipfile.ZipFile(z, "a", zipfile.ZIP_DEFLATED, False) as zf:
        for i, p in enumerate(r.pages):
            w = PdfWriter(); w.add_page(p); b = io.BytesIO(); w.write(b); zf.writestr(f"p_{i+1}.pdf", b.getvalue())
    return z

def merge_pdfs(files):
    w = PdfWriter(); [w.add_page(p) for f in files for p in PdfReader(f).pages]; out = io.BytesIO(); w.write(out); return out

def compress_pdf(file):
    doc = fitz.open(stream=file.getvalue(), filetype="pdf"); out = io.BytesIO(); doc.save(out, deflate=True, garbage=4); return out

def add_watermark(file, text):
    packet = io.BytesIO(); c = canvas.Canvas(packet, pagesize=letter); c.setFont("Helvetica", 50); c.setFillColorRGB(0.5,0.5,0.5,0.3); c.translate(300,400); c.rotate(45); c.drawCentredString(0,0,text); c.save(); packet.seek(0)
    wm = PdfReader(packet).pages[0]; r = PdfReader(file); w = PdfWriter()
    for p in r.pages: p.merge_page(wm); w.add_page(p)
    out = io.BytesIO(); w.write(out); return out

def rotate_pdf(file, angle):
    r = PdfReader(file); w = PdfWriter()
    for p in r.pages: p.rotate(angle); w.add_page(p)
    out = io.BytesIO(); w.write(out); return out

def delete_page(file, p_num):
    r = PdfReader(file); w = PdfWriter()
    for i, p in enumerate(r.pages):
        if i != (p_num-1): w.add_page(p)
    out = io.BytesIO(); w.write(out); return out

def run_ml_pipeline(text):
    try: doc_type = ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': f"Classify into [Resume, Contract, Invoice, Report]. Only name. Text: {text[:1000]}"}])['message']['content'].strip()
    except: doc_type = "Document"
    try: summary = ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': f"3-sentence summary. Text: {text[:4000]}"}])['message']['content']
    except: summary = "Unavailable."
    try: insights = ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': f"Extract key points. Bullets. Text: {text[:4000]}"}])['message']['content']
    except: insights = "Unavailable."
    return doc_type, summary, insights

def chat_with_bookshelf(query, bookshelf):
    context = ""
    for book in bookshelf[:3]:
        try:
            with pdfplumber.open(io.BytesIO(book['data'])) as pdf:
                context += f"\n--- {book['name']} ---\n{''.join([p.extract_text() or '' for p in pdf.pages])[:2000]}..."
        except: pass
    try: return ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': f"Context: {context}\nQ: {query}"}])['message']['content']
    except: return "Error."

def pdf_page_to_image(file_bytes, page_num=0):
    doc = fitz.open(stream=file_bytes, filetype="pdf"); page = doc.load_page(page_num); return page.get_pixmap().tobytes("png")

def analyze_image_with_vision(image_bytes):
    try: return ollama.chat(model='llama3.2-vision', messages=[{'role': 'user', 'content': 'Describe this image.', 'images': [image_bytes]}])['message']['content']
    except: return "Vision Model Missing."

def compare_pdfs_enhanced(text1, text2):
    diff = list(difflib.Differ().compare(text1.splitlines(), text2.splitlines()))
    structured = []; raw = []
    for l in diff:
        if l.startswith('+ '): structured.append(('add', l[2:])); raw.append(f"ADDED: {l[2:]}")
        elif l.startswith('- '): structured.append(('rem', l[2:])); raw.append(f"REMOVED: {l[2:]}")
    return structured, "\n".join(raw)

# --- 7. PREMIUM SIDEBAR NAVIGATION ---
with st.sidebar:
    # Enhanced Header
    st.markdown("""
    <div style='text-align: center; padding: 20px 0; margin-bottom: 20px;'>
        <div style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                    width: 60px; height: 60px; border-radius: 16px; 
                    display: flex; align-items: center; justify-content: center; 
                    margin: 0 auto 12px; box-shadow: 0 8px 24px rgba(102, 126, 234, 0.4);'>
            <span style='font-size: 32px;'>📜</span>
        </div>
        <h2 style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                   -webkit-background-clip: text; -webkit-text-fill-color: transparent; 
                   margin: 0; font-size: 1.5rem; font-weight: 800; font-family: "Poppins", sans-serif;'>
            Journey Through Pages
        </h2>
        <p style='color: rgba(255, 255, 255, 0.7); font-size: 0.75rem; margin-top: 4px; 
                  background: rgba(255, 255, 255, 0.1); padding: 4px 12px; 
                  border-radius: 12px; display: inline-block;'>
            v2.0 Premium
        </p>
    </div>
    """, unsafe_allow_html=True)
    
    ui_spacer(10)
    
    # Enhanced Navigation Menu
    selected = option_menu(
        menu_title="", 
        options=["Home", "Viewer", "Converter", "Editor", "AI Analyst", "Compare", "About"], 
        icons=["house-fill", "eye-fill", "arrow-repeat", "tools", "stars", "arrow-left-right", "info-circle"], 
        default_index=st.session_state.page_selection, 
        manual_select=st.session_state.page_selection,
        styles={
            "container": {
                "padding": "8px 0",
                "background-color": "rgba(255, 255, 255, 0.05)",
                "backdrop-filter": "blur(20px)",
                "border-radius": "16px",
                "border": "1px solid rgba(255, 255, 255, 0.1)"
            },
            "nav-link": {
                "font-size": "0.95rem",
                "text-align": "left",
                "margin": "4px 8px",
                "padding": "12px 16px",
                "border-radius": "12px",
                "transition": "all 0.3s",
                "color": "rgba(255, 255, 255, 0.8)",
                "font-weight": "500"
            },
            "nav-link:hover": {
                "background-color": "rgba(255, 255, 255, 0.1)",
                "transform": "translateX(4px)"
            },
            "nav-link-selected": {
                "background": "linear-gradient(135deg, #667eea 0%, #764ba2 100%)",
                "color": "white",
                "font-weight": "700",
                "box-shadow": "0 4px 14px rgba(102, 126, 234, 0.4)",
                "border": "none"
            },
            "icon": {
                "font-size": "1.2rem",
                "margin-right": "8px"
            }
        }
    )
    
    st.markdown("<div style='height: 20px;'></div>", unsafe_allow_html=True)
    
    # Enhanced Bookshelf Section
    st.markdown("""
    <div style='background: rgba(255, 255, 255, 0.08); backdrop-filter: blur(20px); 
                border-radius: 16px; padding: 16px; border: 1px solid rgba(255, 255, 255, 0.15);'>
        <h4 style='color: rgba(255, 255, 255, 0.9); margin-bottom: 12px; font-size: 1rem; 
                   display: flex; align-items: center; gap: 8px;'>
            <span style='font-size: 1.2rem;'>📚</span> Active Bookshelf
        </h4>
    </div>
    """, unsafe_allow_html=True)
    
    if st.session_state.bookshelf:
        st.markdown("<div style='margin-top: 12px;'>", unsafe_allow_html=True)
        for i, b in enumerate(st.session_state.bookshelf):
            c1, c2 = st.columns([4, 1])
            with c1: 
                st.markdown(f"""
                <div class='book-item' style='margin-bottom: 8px;'>
                    <span style='display: flex; align-items: center; gap: 8px;'>
                        <span style='font-size: 1.1rem;'>📄</span>
                        <span style='font-weight: 500;'>{b['name'][:20]}{'...' if len(b['name']) > 20 else ''}</span>
                    </span>
                </div>
                """, unsafe_allow_html=True)
            with c2: 
                if st.button("✖", key=f"d_{i}", help="Remove from shelf"):
                    remove_from_shelf(b['name'])
                    st.rerun()
        st.markdown("</div>", unsafe_allow_html=True)
    else: 
        st.markdown("""
        <div style='text-align: center; padding: 20px; color: rgba(255, 255, 255, 0.6); 
                    font-size: 0.85rem; font-style: italic;'>
            📖 Library is empty
        </div>
        """, unsafe_allow_html=True)

# --- 8. PAGE LOGIC (WITH REFINED UNIFORM UI) ---
st.markdown("<div class='animate-enter'>", unsafe_allow_html=True)

# Top Navigation Bar
st.markdown("""
<div class="top-nav" style="margin: -1rem -1rem 2rem -1rem; padding: 16px 24px;">
    <div style="text-align: center; width: 100%;">
        <h1 style="font-size: 2.5rem; font-weight: 900; margin: 0; 
                   background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                   -webkit-background-clip: text; -webkit-text-fill-color: transparent; 
                   background-clip: text; letter-spacing: -1px;">
            Journey Through Pages
        </h1>
    </div>
</div>
""", unsafe_allow_html=True)

nav_cols = st.columns(7)
page_buttons = [
    ("🏠", "Home", 0),
    ("👀", "Viewer", 1),
    ("🔄", "Converter", 2),
    ("✂️", "Editor", 3),
    ("🧠", "AI Analyst", 4),
    ("⚖️", "Compare", 5),
    ("ℹ️", "About", 6)
]

for i, (icon, name, idx) in enumerate(page_buttons):
    with nav_cols[i]:
        button_type = "primary" if selected == name else "secondary"
        if st.button(f"{icon} {name}", use_container_width=True, type=button_type, key=f"nav_{name}"):
            set_page(idx)
            st.rerun()

if selected == "Home":
    # --- ENHANCED HERO ---
    st.markdown(f"""
    <div style='text-align:center; padding: 50px 0 40px;'>
        <h1 style='margin-bottom: 24px; font-size: 4rem; font-weight: 900; letter-spacing: -2px;'>
            Journey Through Pages
        </h1>
        <p style='color:#64748b; font-style:italic; font-size: 1.2rem; margin: 0; line-height: 1.8; max-width: 700px; margin: 0 auto;'>
            {get_daily_quote()}
        </p>
    </div>
    """, unsafe_allow_html=True)
    ui_spacer(30)

    # --- FEATURE CARDS (ENHANCED WITH MORE INFO) ---
    r1c1, r1c2, r1c3 = st.columns(3)

    # 👀 VIEWER
    with r1c1:
        st.markdown("""
        <div class="glass-card">
            <div class="home-card-content">
                <h3 style="display: flex; align-items: center; gap: 10px; margin-bottom: 12px;">
                    <span style="font-size: 2rem;">👀</span>
                    <span>Viewer</span>
                </h3>
                <p style="color: #64748b; font-size: 0.9rem; margin-bottom: 16px; line-height: 1.6;">
                    View, annotate, and interact with your PDF documents. Perfect for reviewing and marking up documents.
                </p>
                <ul style="margin-top: auto;">
                    <li><strong>Highlight Text</strong> - Mark important sections with colored highlights</li>
                    <li><strong>Redaction & Shield</strong> - Permanently remove sensitive information</li>
                    <li><strong>Semantic Search</strong> - AI-powered contextual search within documents</li>
                    <li><strong>PII Detection</strong> - Automatically find and protect personal information</li>
                </ul>
            </div>
        </div>
        """, unsafe_allow_html=True)
        if st.button("Open Viewer", key="btn_viewer"): set_page(1); st.rerun()

    # 🔄 CONVERTER
    with r1c2:
        st.markdown("""
        <div class="glass-card">
            <div class="home-card-content">
                <h3 style="display: flex; align-items: center; gap: 10px; margin-bottom: 12px;">
                    <span style="font-size: 2rem;">🔄</span>
                    <span>Converter</span>
                </h3>
                <p style="color: #64748b; font-size: 0.9rem; margin-bottom: 16px; line-height: 1.6;">
                    Convert files between different formats seamlessly. Support for images, Word, Excel, and more.
                </p>
                <ul style="margin-top: auto;">
                    <li><strong>Any-to-PDF</strong> - Convert images, Word docs to PDF</li>
                    <li><strong>PDF-to-Word/Excel</strong> - Extract content to editable formats</li>
                    <li><strong>Image Conversion</strong> - Combine multiple images into PDF</li>
                    <li><strong>Data Export</strong> - Convert CSV/Excel to formatted PDF</li>
                </ul>
            </div>
        </div>
        """, unsafe_allow_html=True)
        if st.button("Open Converter", key="btn_converter"): set_page(2); st.rerun()

    # ✂️ EDITOR
    with r1c3:
        st.markdown("""
        <div class="glass-card">
            <div class="home-card-content">
                <h3 style="display: flex; align-items: center; gap: 10px; margin-bottom: 12px;">
                    <span style="font-size: 2rem;">✂️</span>
                    <span>Editor</span>
                </h3>
                <p style="color: #64748b; font-size: 0.9rem; margin-bottom: 16px; line-height: 1.6;">
                    Professional PDF editing tools to manipulate, organize, and enhance your documents.
                </p>
                <ul style="margin-top: auto;">
                    <li><strong>Split & Merge</strong> - Divide or combine PDF files</li>
                    <li><strong>Compress & Rotate</strong> - Optimize file size and orientation</li>
                    <li><strong>Watermarking</strong> - Add custom watermarks to documents</li>
                    <li><strong>Page Management</strong> - Delete or reorganize pages</li>
                </ul>
            </div>
        </div>
        """, unsafe_allow_html=True)
        if st.button("Open Editor", key="btn_editor"): set_page(3); st.rerun()

    ui_spacer(20)

    # --- ROW 2 ---
    r2c1, r2c2, r2c3 = st.columns(3)

    with r2c1:
        st.markdown("""
        <div class="glass-card">
            <div class="home-card-content">
                <h3 style="display: flex; align-items: center; gap: 10px; margin-bottom: 12px;">
                    <span style="font-size: 2rem;">🧠</span>
                    <span>AI Analyst</span>
                </h3>
                <p style="color: #64748b; font-size: 0.9rem; margin-bottom: 16px; line-height: 1.6;">
                    Advanced AI-powered analysis tools for deep document insights and intelligent processing.
                </p>
                <ul style="margin-top: auto;">
                    <li><strong>AI Summary & Vision</strong> - Intelligent document analysis</li>
                    <li><strong>Audit & Risk Assessment</strong> - Compliance and risk scoring</li>
                    <li><strong>Knowledge Graph</strong> - Visual relationship mapping</li>
                    <li><strong>Timeline & Flashcards</strong> - Study tools and chronologies</li>
                </ul>
            </div>
        </div>
        """, unsafe_allow_html=True)
        if st.button("Open AI Analyst", key="btn_ai"): set_page(4); st.rerun()

    with r2c2:
        st.markdown("""
        <div class="glass-card">
            <div class="home-card-content">
                <h3 style="display: flex; align-items: center; gap: 10px; margin-bottom: 12px;">
                    <span style="font-size: 2rem;">⚖️</span>
                    <span>Compare</span>
                </h3>
                <p style="color: #64748b; font-size: 0.9rem; margin-bottom: 16px; line-height: 1.6;">
                    Compare two versions of documents to identify changes, additions, and deletions.
                </p>
                <ul style="margin-top: auto;">
                    <li><strong>Version Comparison</strong> - Side-by-side document diff</li>
                    <li><strong>Change Detection</strong> - Highlight additions and removals</li>
                    <li><strong>AI Summary</strong> - Intelligent change summaries</li>
                    <li><strong>Visual Diff</strong> - Color-coded change indicators</li>
                </ul>
            </div>
        </div>
        """, unsafe_allow_html=True)
        if st.button("Open Compare", key="btn_compare"): set_page(5); st.rerun()

    with r2c3:
        st.markdown("""
        <div class="glass-card">
            <div class="home-card-content">
                <h3 style="display: flex; align-items: center; gap: 10px; margin-bottom: 12px;">
                    <span style="font-size: 2rem;">📚</span>
                    <span>Book Shelf</span>
                </h3>
                <p style="color: #64748b; font-size: 0.9rem; margin-bottom: 16px; line-height: 1.6;">
                    Your personal document library. Securely stored for quick access and AI-powered cross-document queries.
                </p>
                <ul style="margin-top: auto;">
                    <li><strong>Quick Access</strong> - Recently used documents</li>
                    <li><strong>Secure Storage</strong> - Privacy-first document management</li>
                    <li><strong>AI Queries</strong> - Cross-document intelligent search</li>
                    <li><strong>Omniscient Mode</strong> - Context-aware AI responses</li>
                </ul>
            </div>
        </div>
        """, unsafe_allow_html=True)
        

elif selected == "Viewer":
    st.markdown("<h1>👀 Ethereal Viewer</h1>", unsafe_allow_html=True)
    c_tools, c_view = st.columns([1, 2])
    
    with c_tools:
        # Enhanced Container for Inputs
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 20px;'>
                <h3 style='display: flex; align-items: center; gap: 10px; margin-bottom: 16px;'>
                    <span style='font-size: 1.5rem;'>📂</span> Source Document
                </h3>
            </div>
            """, unsafe_allow_html=True)
            bk = [b['name'] for b in st.session_state.bookshelf]
            sel = st.selectbox(
                "📚 Select from Library", 
                ["None"] + bk,
                help="Choose a document from your bookshelf"
            )
            st.markdown("""
            <div style='text-align: center; padding: 12px 0; color: rgba(102, 126, 234, 0.7); 
                        font-weight: 500; font-size: 0.9rem;'>
                ─── OR ───
            </div>
            """, unsafe_allow_html=True)
            up = st.file_uploader(
                "📤 Upload New PDF", 
                type="pdf",
                help="Upload a PDF file to view and edit"
            )
        
        target = None; name = None
        if up: target = up.getvalue(); name = up.name
        elif sel != "None": target = get_file_from_shelf(sel).getvalue(); name = sel
        
        if target:
            if "curr_view" not in st.session_state or st.session_state.curr_view != name:
                st.session_state.curr_view = name; st.session_state.highlights = []; st.session_state.view_bytes = target; st.session_state.orig_bytes = target; add_to_bookshelf(name, target)
            
            ui_spacer(10)
            t1, t2, t3, t4 = st.tabs(["✨ Highlight", "🔒 Redact", "🛡️ Shield", "🔍 Search"])
            
            with t1:
                st.markdown("""
                <div style='margin-bottom: 12px;'>
                    <p style='color: #64748b; font-size: 0.9rem;'>✨ Mark important sections with highlights</p>
                </div>
                """, unsafe_allow_html=True)
                hl = st.text_input("📝 Text to Mark", placeholder="Enter text to highlight...", help="Enter the exact text you want to highlight")
                c1, c2 = st.columns(2)
                if c1.button("Highlight"):
                    st.session_state.highlights.append({'text': hl})
                    res, _ = apply_highlights_to_pdf(st.session_state.orig_bytes, st.session_state.highlights)
                    st.session_state.view_bytes = res.getvalue(); st.rerun()
                if c2.button("Reset"): st.session_state.highlights = []; st.session_state.view_bytes = st.session_state.orig_bytes; st.rerun()
            with t2:
                st.markdown("""
                <div style='margin-bottom: 12px;'>
                    <p style='color: #64748b; font-size: 0.9rem;'>🔒 Permanently remove sensitive text from document</p>
                </div>
                """, unsafe_allow_html=True)
                rd = st.text_input("🔒 Text to Redact", placeholder="Enter confidential information...", help="Enter text that should be permanently blacked out")
                if st.button("Apply Redaction"):
                    res, c = apply_redaction(st.session_state.view_bytes, rd)
                    st.session_state.view_bytes = res.getvalue(); st.session_state.orig_bytes = res.getvalue(); st.success(f"Redacted {c}."); st.rerun()
            with t3:
                st.markdown("""
                <div style='margin-bottom: 12px;'>
                    <p style='color: #64748b; font-size: 0.9rem;'>🛡️ Automatically detect and redact PII (Emails, Phone Numbers)</p>
                </div>
                """, unsafe_allow_html=True)
                if st.button("🛡️ Auto Scan & Redact", use_container_width=True):
                    txt = extract_text_with_references(io.BytesIO(st.session_state.view_bytes)); ft = " ".join([p['text'] for p in txt])
                    pii = detect_pii(ft)
                    if pii:
                        tb = st.session_state.view_bytes; tot = 0
                        for i in pii: rb, c = apply_redaction(tb, i); tb = rb.getvalue(); tot += c
                        st.session_state.view_bytes = tb; st.session_state.orig_bytes = tb; st.success(f"Redacted {tot} items."); st.rerun()
                    else: st.info("No PII found.")
            with t4:
                st.markdown("""
                <div style='margin-bottom: 12px;'>
                    <p style='color: #64748b; font-size: 0.9rem;'>🔍 Find contextually similar text using AI semantic search</p>
                </div>
                """, unsafe_allow_html=True)
                sq = st.text_input("🔍 Semantic Query", placeholder="What is the invoice total?", help="Ask a question or describe what you're looking for")
                if st.button("🧠 Deep Search"):
                    sk = ui_skeleton()
                    res, sc = semantic_search_pdf(st.session_state.view_bytes, sq)
                    sk.empty()
                    if res and sc > 0.3: 
                        st.session_state.view_bytes = res.getvalue(); ui_confidence_badge("High"); st.rerun()
                    else: st.error("No matches found.")
            
            ui_spacer(20)
            st.download_button("💾 Download Result", st.session_state.view_bytes, "doc.pdf", "application/pdf", use_container_width=True)

    with c_view:
        if target:
            # Ensure view_bytes is initialized
            if "view_bytes" not in st.session_state:
                st.session_state.view_bytes = target
                st.session_state.orig_bytes = target
            
            # Display the PDF
            try:
                pdf_bytes = st.session_state.view_bytes if "view_bytes" in st.session_state else target
                b64 = base64.b64encode(pdf_bytes).decode('utf-8')
                st.markdown(f'<iframe src="data:application/pdf;base64,{b64}#t={time.time()}" width="100%" height="850" style="border-radius:15px; border:1px solid #ddd; background:white;"></iframe>', unsafe_allow_html=True)
            except Exception as e:
                st.error(f"Error displaying PDF: {str(e)}")
                # Fallback: try to display directly from target
                try:
                    b64 = base64.b64encode(target).decode('utf-8')
                    st.markdown(f'<iframe src="data:application/pdf;base64,{b64}#t={time.time()}" width="100%" height="850" style="border-radius:15px; border:1px solid #ddd; background:white;"></iframe>', unsafe_allow_html=True)
                    st.session_state.view_bytes = target
                    st.session_state.orig_bytes = target
                except:
                    st.error("Unable to display PDF. Please try uploading again.")
        else:
            st.info("👈 Select a document to view")

elif selected == "AI Analyst":
    st.markdown("<h1>🧠 The AI Analyst</h1>", unsafe_allow_html=True)
    
    # Enhanced Tabs
    tabs = st.tabs(["📄 Text Analysis", "👁️ Vision", "🎓 Study Tools", "🕸️ Knowledge Graph", "⚖️ Auditor", "🎞️ Presentation", "📊 Ledger", "⏳ Timeline", "🕵️ Truth Check"])
    
    with tabs[0]:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>📄</span> Document Analysis
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Upload a PDF for AI-powered text analysis and insights</p>
            </div>
            """, unsafe_allow_html=True)
            up = st.file_uploader("Choose PDF file", type="pdf", key="ai", help="Upload a document for AI analysis")
        
        if up:
            if "curr_ai" not in st.session_state or st.session_state.curr_ai != up.name:
                st.session_state.curr_ai = up.name; sk = ui_skeleton()
                pgs = extract_text_with_references(up)
                ft = "\n".join([p['text'] for p in pgs])
                dt, sm, ins = run_ml_pipeline(ft)
                st.session_state.dt = dt; st.session_state.sm = sm; st.session_state.ins = ins; st.session_state.ft = ft
                sk.empty(); st.rerun()
            
            col_a, col_b = st.columns([2, 1])
            with col_a:
                st.subheader("Summary")
                st.info(st.session_state.sm)
            with col_b:
                st.subheader("Metadata")
                st.write(f"**Type:** {st.session_state.dt}")
                if st.button("🎧 Play Audio Summary"):
                    with st.spinner("Synthesizing..."):
                        ad = text_to_audio_mac(st.session_state.ft[:1000])
                        if ad: st.audio(ad, format="audio/wav")
            
            st.divider()
            st.subheader("💬 Chat with Document")
            use_shelf = st.checkbox("📚 Omniscient Mode (Include Bookshelf Context)")
            if p := st.chat_input("Ask about this document..."):
                if use_shelf:
                    sk = ui_skeleton(); ans = chat_with_bookshelf(p, st.session_state.bookshelf); sk.empty(); st.write(f"**AI:** {ans}")
                else:
                    sk = ui_skeleton(); res = ollama.chat(model='llama3.2', messages=[{'role':'user','content':f"Context: {st.session_state.ft[:10000]} Q: {p}"}]); sk.empty()
                    st.write(f"**AI:** {res['message']['content']}")

    with tabs[1]:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>👁️</span> Vision Analysis
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Analyze visual content in PDF pages using AI vision</p>
            </div>
            """, unsafe_allow_html=True)
            v_up = st.file_uploader("Choose PDF with images", type="pdf", key="vis", help="Upload a PDF containing images to analyze")
        if v_up:
            c1, c2 = st.columns(2)
            img = pdf_page_to_image(v_up.getvalue(), 0)
            c1.image(img, caption="First Page Preview", use_column_width=True)
            with c2:
                if st.button("Analyze Visual Content", use_container_width=True):
                    sk = ui_skeleton(); desc = analyze_image_with_vision(img); sk.empty()
                    st.success("Analysis Complete"); st.write(desc)

    with tabs[2]:
        if "ft" in st.session_state:
            st.caption("Generate study materials from the uploaded doc")
            if st.button("Generate Flashcards", use_container_width=True):
                sk = ui_skeleton(); cards = generate_flashcards(st.session_state.ft); sk.empty()
                if cards: st.table(pd.DataFrame(cards)); ui_confidence_badge("High")
        else: st.warning("Please upload a document in the 'Text' tab first.")

    with tabs[3]:
        if "ft" in st.session_state:
            if st.button("Build Knowledge Graph", use_container_width=True):
                nlp = spacy.load("en_core_web_sm"); doc = nlp(st.session_state.ft[:5000])
                nodes = []; edges = []; names = set()
                for e in doc.ents:
                    if e.label_ in ["ORG","PERSON"] and e.text not in names:
                        nodes.append(Node(id=e.text, label=e.text, size=20)); names.add(e.text)
                for s in doc.sents:
                    en = [x.text for x in s.ents if x.text in names]
                    for i in range(len(en)-1): edges.append(Edge(source=en[i], target=en[i+1]))
                agraph(nodes, edges, Config(height=500, width=700))
        else: st.warning("Upload document first.")

    with tabs[4]:
        if "ft" in st.session_state:
            if st.button("Run Compliance Audit", use_container_width=True):
                sk = ui_skeleton(); res = run_audit(st.session_state.ft); sk.empty()
                st.markdown(f"<div class='glass-card'>{res}</div>", unsafe_allow_html=True)

    with tabs[5]:
        if "ft" in st.session_state:
            if st.button("Generate PowerPoint Deck", use_container_width=True):
                sk = ui_skeleton(); ppt = generate_ppt_from_text(st.session_state.ft); sk.empty()
                if ppt: st.download_button("Download PPT", ppt.getvalue(), "pres.pptx", use_container_width=True); ui_confidence_badge("High")

    with tabs[6]:
        if "ft" in st.session_state:
            if st.button("Extract Financial Data", use_container_width=True):
                sk = ui_skeleton(); data = extract_ledger_data(st.session_state.ft); sk.empty()
                if data: st.json(data); ui_confidence_badge("High")

    with tabs[7]:
        if "ft" in st.session_state:
            if st.button("Build Chronological Timeline", use_container_width=True):
                sk = ui_skeleton(); data = extract_timeline_data(st.session_state.ft); sk.empty()
                if data and "events" in data:
                    timeline(data, height=400)
                    ui_confidence_badge("Medium")
                else: st.warning("No timeline data found.")

    with tabs[8]:
        if "ft" in st.session_state:
            if st.button("Analyze Logic & Contradictions", use_container_width=True):
                sk = ui_skeleton(); res = analyze_contradictions(st.session_state.ft); sk.empty()
                st.markdown(f"<div class='glass-card'>{res}</div>", unsafe_allow_html=True)

elif selected == "Converter":
    st.markdown("<h1>🔄 Universal Converter</h1>", unsafe_allow_html=True)
    t1, t2, t3, t4, t5 = st.tabs(["🖼️ Images to PDF", "📝 Word to PDF", "📄 PDF to Word", "📊 PDF to Excel", "💾 Data to PDF"])
    
    with t1:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>🖼️</span> Upload Images
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Select multiple PNG or JPG images to combine into a single PDF</p>
            </div>
            """, unsafe_allow_html=True)
            i = st.file_uploader(
                "Choose image files", 
                type=["png","jpg","jpeg"], 
                accept_multiple_files=True,
                help="You can select multiple images at once"
            )
            if i:
                st.success(f"✅ {len(i)} image(s) selected")
                ui_spacer(10)
                if st.button("🔄 Convert to PDF", key="btn_i2p", use_container_width=True): 
                    st.download_button("💾 Download PDF", images_to_pdf(i).getvalue(), "img.pdf", use_container_width=True)
    with t2:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>📝</span> Upload Word Document
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Convert your DOCX file to PDF format</p>
            </div>
            """, unsafe_allow_html=True)
            d = st.file_uploader("Choose DOCX file", type=["docx"], help="Upload a Microsoft Word document")
            if d:
                st.success(f"✅ {d.name} selected")
                ui_spacer(10)
                if st.button("🔄 Convert to PDF", key="btn_w2p", use_container_width=True): 
                    st.download_button("💾 Download PDF", word_to_pdf(d).getvalue(), "doc.pdf", use_container_width=True)
    with t3:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>📄</span> Upload PDF
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Convert PDF to editable Word document</p>
            </div>
            """, unsafe_allow_html=True)
            p = st.file_uploader("Choose PDF file", type="pdf", key="p2w", help="Upload a PDF file to convert")
            if p:
                st.success(f"✅ {p.name} selected")
                ui_spacer(10)
                if st.button("🔄 Convert to Word", use_container_width=True): 
                    sk = ui_skeleton()
                    doc = pdf_to_docx(p)
                    sk.empty()
                    st.download_button("💾 Download Docx", doc.getvalue(), "doc.docx", use_container_width=True)
    with t4:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>📊</span> Upload PDF
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Extract tables and data from PDF to Excel</p>
            </div>
            """, unsafe_allow_html=True)
            p2 = st.file_uploader("Choose PDF file", type="pdf", key="p2e", help="Upload a PDF with tables to extract")
            if p2:
                st.success(f"✅ {p2.name} selected")
                ui_spacer(10)
                if st.button("🔄 Convert to Excel", use_container_width=True):
                    sk = ui_skeleton()
                    xl = pdf_to_excel(p2)
                    sk.empty()
                    if xl: st.download_button("💾 Download Excel", xl.getvalue(), "tab.xlsx", use_container_width=True)
    with t5:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>💾</span> Upload Data File
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Convert CSV or Excel files to formatted PDF</p>
            </div>
            """, unsafe_allow_html=True)
            d = st.file_uploader("Choose CSV/Excel file", type=["csv","xlsx"], help="Upload CSV or Excel file")
            if d:
                st.success(f"✅ {d.name} selected")
                ui_spacer(10)
                if st.button("🔄 Convert to PDF", key="btn_d2p", use_container_width=True):
                    sk = ui_skeleton()
                    if d.name.endswith("csv"): df = pd.read_csv(d)
                    else: df = pd.read_excel(d)
                    pdf = df_to_pdf(df)
                    sk.empty()
                    st.download_button("💾 Download PDF", pdf.getvalue(), "data.pdf", use_container_width=True)

elif selected == "Editor":
    st.markdown("<h1>✂️ Pro Editor</h1>", unsafe_allow_html=True)
    t1, t2, t3, t4, t5, t6 = st.tabs(["✂️ Split", "🔗 Merge", "🗜️ Compress", "💧 Watermark", "🔄 Rotate", "🗑️ Delete"])
    
    with t1:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>✂️</span> Split PDF into Individual Pages
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Separate each page into its own PDF file</p>
            </div>
            """, unsafe_allow_html=True)
            f = st.file_uploader("Choose PDF file", type="pdf", key="sp", help="Select a PDF to split into pages")
            if f:
                st.success(f"✅ {f.name} selected")
                if st.button("✂️ Split Pages", use_container_width=True): 
                    r = PdfReader(f)
                    z = io.BytesIO()
                with zipfile.ZipFile(z,"a") as zf:
                    for n, p in enumerate(r.pages):
                            w = PdfWriter()
                            w.add_page(p)
                            b = io.BytesIO()
                            w.write(b)
                            zf.writestr(f"page_{n+1}.pdf", b.getvalue())
                    st.download_button("💾 Download ZIP", z.getvalue(), "pages.zip", use_container_width=True)
    with t2:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>🔗</span> Merge Multiple PDFs
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Combine multiple PDF files into one document</p>
            </div>
            """, unsafe_allow_html=True)
            fs = st.file_uploader("Choose PDF files", accept_multiple_files=True, key="mg", help="Select multiple PDFs to merge")
            if fs:
                st.success(f"✅ {len(fs)} file(s) selected")
                if st.button("🔗 Merge Files", use_container_width=True): 
                    st.download_button("💾 Download Merged PDF", merge_pdfs(fs).getvalue(), "merged.pdf", use_container_width=True)
    with t3:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>🗜️</span> Compress PDF
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Reduce PDF file size while maintaining quality</p>
            </div>
            """, unsafe_allow_html=True)
            fc = st.file_uploader("Choose PDF file", type="pdf", key="cp", help="Select a PDF to compress")
            if fc:
                st.success(f"✅ {fc.name} selected")
                if st.button("🗜️ Compress File", use_container_width=True): 
                    st.download_button("💾 Download Compressed PDF", compress_pdf(fc).getvalue(), "compressed.pdf", use_container_width=True)
    with t4:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>💧</span> Add Watermark
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Overlay text watermark on your PDF</p>
            </div>
            """, unsafe_allow_html=True)
            fw = st.file_uploader("Choose PDF file", type="pdf", key="wm", help="Select a PDF to watermark")
            txt = st.text_input("💬 Watermark Text", "DRAFT", help="Enter the text to display as watermark")
            if fw:
                st.success(f"✅ {fw.name} selected")
                if st.button("💧 Apply Watermark", use_container_width=True): 
                    st.download_button("💾 Download PDF", add_watermark(fw, txt).getvalue(), "watermarked.pdf", use_container_width=True)
    with t5:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>🔄</span> Rotate Pages
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Rotate all pages in your PDF document</p>
            </div>
            """, unsafe_allow_html=True)
            fr = st.file_uploader("Choose PDF file", type="pdf", key="rt", help="Select a PDF to rotate")
            ang = st.radio("📐 Rotation Angle", [90, 180, 270], horizontal=True, help="Choose the rotation angle")
            if fr:
                st.success(f"✅ {fr.name} selected")
                if st.button("🔄 Rotate Pages", use_container_width=True): 
                    st.download_button("💾 Download PDF", rotate_pdf(fr, ang).getvalue(), "rotated.pdf", use_container_width=True)
    with t6:
        with st.container(border=True):
            st.markdown("""
            <div style='margin-bottom: 16px;'>
                <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                    <span>🗑️</span> Delete Page
                </h4>
                <p style='color: #64748b; font-size: 0.9rem;'>Remove a specific page from your PDF</p>
            </div>
            """, unsafe_allow_html=True)
            fd = st.file_uploader("Choose PDF file", type="pdf", key="del", help="Select a PDF to edit")
            if fd:
                r = PdfReader(fd)
                total_pages = len(r.pages)
                pn = st.number_input(
                    f"📄 Page Number to Delete (1-{total_pages})", 
                    min_value=1, 
                    max_value=total_pages, 
                    value=1,
                    help=f"Select which page to remove (total pages: {total_pages})"
                )
                st.success(f"✅ {fd.name} selected ({total_pages} pages)")
                if st.button("🗑️ Delete Page", use_container_width=True): 
                    st.download_button("💾 Download PDF", delete_page(fd, pn).getvalue(), "modified.pdf", use_container_width=True)

elif selected == "Compare":
    st.markdown("<h1>⚖️ Cross-Comparison</h1>", unsafe_allow_html=True)
    
    with st.container(border=True):
        st.markdown("""
        <div style='margin-bottom: 20px;'>
            <h4 style='display: flex; align-items: center; gap: 8px; margin-bottom: 8px;'>
                <span>⚖️</span> Document Comparison
            </h4>
            <p style='color: #64748b; font-size: 0.9rem;'>Compare two PDF versions to identify differences and changes</p>
        </div>
        """, unsafe_allow_html=True)
        c1, c2 = st.columns(2)
        with c1: 
            f1 = st.file_uploader(
                "📄 Original Version", 
                type="pdf", 
                key="c1",
                help="Upload the original PDF document"
            )
        with c2: 
            f2 = st.file_uploader(
                "📝 Modified Version", 
                type="pdf", 
                key="c2",
                help="Upload the modified PDF document"
            )
    
    if f1 and f2:
        if st.button("Run Comparison Analysis", use_container_width=True):
            sk = ui_skeleton()
            t1 = "".join([p.extract_text() for p in PdfReader(f1).pages])
            t2 = "".join([p.extract_text() for p in PdfReader(f2).pages])
            d, raw = compare_pdfs_enhanced(t1, t2)
            
            st.subheader("📝 AI Change Summary")
            try: 
                s = ollama.chat(model='llama3.2', messages=[{'role': 'user', 'content': f"Summarize changes:\n{raw[:4000]}"}])['message']['content']
                st.info(s)
            except: st.warning("Summary unavailable.")
            
            sk.empty()
            st.subheader("Detailed Diff")
            for t, txt in d:
                if t == 'add': st.markdown(f"<div class='diff-add'>{txt}</div>", unsafe_allow_html=True)
                elif t == 'rem': st.markdown(f"<div class='diff-rem'>{txt}</div>", unsafe_allow_html=True)

elif selected == "About":
    st.markdown("""
    <div style='text-align: center; padding: 30px 0;'>
        <h1>About Journey Through Pages</h1>
    </div>
    """, unsafe_allow_html=True)
    
    # What This Application Does
    st.markdown("""
    <div class="glass-card">
        <h2 style='margin-bottom: 20px; display: flex; align-items: center; gap: 10px;'>
            <span style='font-size: 2rem;'>📖</span>
            <span>What This Application Does</span>
        </h2>
        <p style='color: #475569; font-size: 1.05rem; line-height: 1.8; margin-bottom: 20px;'>
            <strong>Journey Through Pages</strong> is a comprehensive, all-in-one PDF processing platform designed to 
            empower your document workflow. Whether you need to view, convert, edit, analyze, or compare PDFs, 
            this application provides professional-grade tools with a focus on security, privacy, and offline capability.
        </p>
        <div style='display: grid; grid-template-columns: repeat(auto-fit, minmax(250px, 1fr)); gap: 16px; margin-top: 24px;'>
            <div style='padding: 16px; background: rgba(102, 126, 234, 0.1); border-radius: 12px; border-left: 4px solid #667eea;'>
                <h4 style='margin: 0 0 8px 0; color: #667eea;'>🔒 Security First</h4>
                <p style='margin: 0; color: #64748b; font-size: 0.9rem;'>All processing happens locally. Your documents never leave your device.</p>
            </div>
            <div style='padding: 16px; background: rgba(102, 126, 234, 0.1); border-radius: 12px; border-left: 4px solid #667eea;'>
                <h4 style='margin: 0 0 8px 0; color: #667eea;'>💻 Offline Available</h4>
                <p style='margin: 0; color: #64748b; font-size: 0.9rem;'>Works completely offline. No internet required for core features.</p>
            </div>
            <div style='padding: 16px; background: rgba(102, 126, 234, 0.1); border-radius: 12px; border-left: 4px solid #667eea;'>
                <h4 style='margin: 0 0 8px 0; color: #667eea;'>🧠 AI-Powered</h4>
                <p style='margin: 0; color: #64748b; font-size: 0.9rem;'>Advanced AI analysis, semantic search, and intelligent document understanding.</p>
            </div>
            <div style='padding: 16px; background: rgba(102, 126, 234, 0.1); border-radius: 12px; border-left: 4px solid #667eea;'>
                <h4 style='margin: 0 0 8px 0; color: #667eea;'>🆓 Free & Open</h4>
                <p style='margin: 0; color: #64748b; font-size: 0.9rem;'>No subscriptions, no hidden costs. All features available to everyone.</p>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    ui_spacer(30)
    
    # Version Information
    st.markdown("""
    <div class="glass-card">
        <h2 style='margin-bottom: 20px; display: flex; align-items: center; gap: 10px;'>
            <span style='font-size: 2rem;'>📋</span>
            <span>Version 1.2 - What's New</span>
        </h2>
            <h3 style='color: #667eea; margin-top: 0;'>✨ UI Update</h3>
            <ul style='color: #475569; line-height: 1.8; padding-left: 24px;'>
                <li>Complete UI redesign with modern glassmorphism effects</li>
                <li>Enhanced navigation with top bar and improved sidebar</li>
                <li>Better visual hierarchy and consistent design language</li>
                <li>Improved card layouts with detailed feature descriptions</li>
                <li>Enhanced loading indicators with smooth animations</li>
            </ul>
            

    </div>
    """, unsafe_allow_html=True)
    
    ui_spacer(30)
    
    # A Little Story
    st.markdown("""
    <div class="glass-card">
        <h2 style='margin-bottom: 20px; display: flex; align-items: center; gap: 10px;'>
            <span style='font-size: 2rem;'>📖</span>
            <span>A Little Story</span>
        </h2>
        <div style='background: linear-gradient(135deg, rgba(102, 126, 234, 0.05), rgba(118, 75, 162, 0.05)); 
                    padding: 30px; border-radius: 16px; border: 2px solid rgba(102, 126, 234, 0.2);'>
            <p style='color: #475569; font-size: 1.1rem; line-height: 1.9; font-style: italic; margin: 0; text-align: justify;'>
                I have been using many PDF applications but didn't get an all-in-one usable and offline available 
                and security first application. Sometimes I had to pay a huge amount for premium features that should 
                be accessible to everyone. So I had to make this application to power up PDFs - combining all the tools 
                I needed in one place, ensuring privacy and security, and making it available offline so you're never 
                dependent on internet connectivity. This is my solution to the problem of fragmented, expensive, and 
                privacy-concerning PDF tools.
            </p>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    ui_spacer(30)
    
    # Made By
    st.markdown("""
    <div class="glass-card" style='text-align: center;'>
        <div style='padding: 30px 20px;'>
            <div style='font-size: 2.5rem; margin-bottom: 16px;'>👨‍💻</div>
            <h2 style='margin-bottom: 10px; color: #1e293b; font-size: 1.2rem;'>Made by</h2>
            <h1 style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                       -webkit-background-clip: text; -webkit-text-fill-color: transparent; 
                       margin: 0; font-size: 1.8rem; font-weight: 800;'>
                M. Ihthisham Irshad
            </h1>
            <p style='color: #64748b; font-size: 0.95rem; margin-top: 12px; margin-bottom: 20px;'>
                Built with ❤️ for everyone who values privacy, security, and powerful tools
            </p>
            <div style='display: flex; justify-content: center; gap: 20px; margin-top: 20px; flex-wrap: wrap;'>
                <a href='https://www.linkedin.com/in/ihthisham-irshad-296157234' target='_blank' 
                   style='display: inline-flex; align-items: center; gap: 8px; padding: 10px 20px; 
                          background: linear-gradient(135deg, #0077b5 0%, #005885 100%); 
                          color: white; text-decoration: none; border-radius: 12px; 
                          font-weight: 600; transition: all 0.3s; box-shadow: 0 4px 12px rgba(0, 119, 181, 0.3);'>
                    <span style='font-size: 1.2rem;'>💼</span>
                    <span>LinkedIn: Ihthisham Irshad</span>
                </a>
                <a href='https://www.instagram.com/inthi.ii?igsh=bW5tZXZ4dmpvYnli&utm_source=qr' target='_blank' 
                   style='display: inline-flex; align-items: center; gap: 8px; padding: 10px 20px; 
                          background: linear-gradient(135deg, #E4405F 0%, #C13584 100%); 
                          color: white; text-decoration: none; border-radius: 12px; 
                          font-weight: 600; transition: all 0.3s; box-shadow: 0 4px 12px rgba(225, 64, 95, 0.3);'>
                    <span style='font-size: 1.2rem;'>📷</span>
                    <span>Instagram: inth.ii</span>
                </a>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)

st.markdown("</div>", unsafe_allow_html=True) # End Animation